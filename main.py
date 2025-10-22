from fastapi import FastAPI, Request, HTTPException
from fastapi.responses import FileResponse
from pathlib import Path
import os

app = FastAPI()

# 创建必要目录
OUTPUT_DIR = Path("output")
TEMPLATES_DIR = Path("templates")
OUTPUT_DIR.mkdir(exist_ok=True)
TEMPLATES_DIR.mkdir(exist_ok=True)

# 模板路径
TEMPLATE_PATH = TEMPLATES_DIR / "template.potx"

@app.post("/generate")
async def generate_ppt(request: Request):
    try:
        # ✅ 正确解析 JSON 请求体
        body = await request.json()
        
        # 确保 body 是 dict 类型
        if not isinstance(body, dict):
            raise HTTPException(status_code=400, detail="请求体必须是 JSON 对象")

        slides = body.get("slides", [])
        if not slides:
            raise HTTPException(status_code=400, detail="缺少幻灯片数据")

        from pptx import Presentation

        # 判断是否使用模板
        if TEMPLATE_PATH.exists():
            prs = Presentation(TEMPLATE_PATH)
        else:
            prs = Presentation()  # 使用默认样式

        # 添加幻灯片
        for slide_data in slides:
            title = slide_data.get("title", "无标题")
            content = slide_data.get("content", "无内容")

            # 使用布局 1：标题和内容
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # 设置标题
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title

            # 设置内容
            for shape in slide.placeholders:
                if hasattr(shape, 'placeholder_format') and \
                   hasattr(shape.placeholder_format, 'idx') and \
                   shape.placeholder_format.idx == 1:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    p.text = content
                    break

        # 保存 PPT 文件
        output_path = OUTPUT_DIR / "output.pptx"
        prs.save(output_path)

        # ✅ 动态生成当前主机的下载链接
        host = f"https://{request.headers['host']}"
        download_link = f"{host}/download/output.pptx"

        # 可选：返回在线预览链接（Office Online）
        preview_url = f"https://view.officeapps.live.com/op/view.aspx?src={download_link}"

        return {
            "success": True,
            "download_link": download_link,
            "preview_url": preview_url
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"生成PPT失败: {str(e)}")


@app.get("/download/{filename}")
async def download_file(filename: str):
    file_path = OUTPUT_DIR / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="文件未找到")
    return FileResponse(
        file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename
    )


# 根路径提示
@app.get("/")
async def home():
    return {
        "message": "PPT 生成服务运行中",
        "endpoints": {
            "生成PPT": "/generate (POST)",
            "下载文件": "/download/{filename} (GET)"
        }
    }